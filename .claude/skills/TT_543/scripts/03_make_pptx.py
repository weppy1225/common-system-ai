#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
[TT_543] 4단계 — 운영자 매뉴얼 PPTX 생성 (템플릿 기반, python-pptx, Windows 경로)

입력:
  deliverables/30-output/05 이행(TT)/tmp_543/capture_config.json
  deliverables/30-output/05 이행(TT)/tmp_543/screens/{메뉴코드}/*.png
  deliverables/30-output/05 이행(TT)/tmp_543/screens/{메뉴코드}/coords.json

템플릿:
  deliverables/10-templates/05 이행(TT)/TT.412. 운영자메뉴얼.pptx
  - 13.33 × 7.5 인치 (16:9)
  - 제목바 #2D4B73 / 이미지영역 0~10in / 설명패널 10~13.33in / 페이지번호 우하단

출력:
  deliverables/30-output/05 이행(TT)/TT_543_운영자매뉴얼_{고객사명}.pptx

설계 원칙 (TT_541 과 동일):
  - 라벨/테두리/배지/커넥터/설명패널은 모두 python-pptx 도형(`add_shape`)으로 PPT 안에
    직접 그린다. 이미지 합성하지 않으므로 PowerPoint 에서 자유롭게 편집 가능.
  - 이미지는 add_picture 의 width/height 직접 지정 (왜곡 방지).
  - 배지는 이미지 우측(IMG_R) ~ 설명 패널 사이 "배지 존" 에만 배치.

TT_541 과의 차이점:
  - 표지 제목 : "사용자 매뉴얼" → "운영자 매뉴얼"
  - 출력 파일명: TT_541_사용자매뉴얼_... → TT_543_운영자매뉴얼_...
  - TMP 디렉토리: tmp → tmp_543
  - 카테고리 분리 표지 (시스템관리 / 마스터 / 권한 등) 선택 지원
  - 설명 패널 디폴트 문구를 운영자 톤으로 살짝 조정
"""
from __future__ import annotations

import json
import os
import re
import sys
import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image, ImageFont

# Vue 소스 파서 (같은 디렉토리)
sys.path.insert(0, str(Path(__file__).resolve().parent))
from parse_vue_source import parse_menu as parse_vue_menu

# ── 경로 (Windows 네이티브 + WSL 호환) ──────────────────────────
# 본 스크립트는 .claude/skills/TT_543/scripts/ 안에 있다.
# parents[0]=scripts  parents[1]=TT_543  parents[2]=skills
# parents[3]=.claude  parents[4]=<프로젝트 루트>
REPO_ROOT = Path(__file__).resolve().parents[4]
TEMPLATE = REPO_ROOT / 'deliverables' / '10-templates' / '05 이행(TT)' / 'TT.412. 운영자메뉴얼.pptx'
OUTPUT_DIR = REPO_ROOT / 'deliverables' / '30-output' / '05 이행(TT)'
TMP_DIR = OUTPUT_DIR / 'tmp_543'
SCREENS_ROOT = TMP_DIR / 'screens'
CFG_FILE = TMP_DIR / 'capture_config.json'

if not CFG_FILE.exists():
    sys.exit(f"[ERR] config 파일이 없습니다: {CFG_FILE}")
if not TEMPLATE.exists():
    sys.exit(f"[ERR] 템플릿 파일이 없습니다: {TEMPLATE}")

cfg = json.loads(CFG_FILE.read_text(encoding='utf-8'))


def sanitize_filename(s: str) -> str:
    return re.sub(r'[\\/:*?"<>|]', '_', str(s)).strip() or 'WMS'


CUSTOMER = sanitize_filename(cfg.get('customer') or 'WMS')
OUT_FILE = OUTPUT_DIR / f"TT_543_운영자매뉴얼_{CUSTOMER}.pptx"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ── 색상 상수 (TT_541 과 동일, 변경 금지) ───────────────────────
COLOR_RED = RGBColor(0xDC, 0x1E, 0x1E)
COLOR_ORANGE = RGBColor(0xC8, 0x6E, 0x00)
COLOR_BLUE = RGBColor(0x1E, 0x64, 0xC8)
COLOR_GREEN = RGBColor(0x14, 0x8C, 0x3C)
COLOR_GRAY = RGBColor(0x6E, 0x6E, 0x6E)
COLOR_NAVY = RGBColor(0x1A, 0x3A, 0x5C)
COLOR_DARK = RGBColor(0x33, 0x33, 0x33)
COLOR_WARN = RGBColor(0xCC, 0x22, 0x22)
COLOR_TITLE_BG = RGBColor(0x2D, 0x4B, 0x73)
COLOR_PANEL_LINE = RGBColor(0xD0, 0xD5, 0xDD)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_PAGE_NUM = RGBColor(0x88, 0x88, 0x88)
COLOR_SUBTITLE = RGBColor(0xD6, 0xE0, 0xF0)

REGION_COLOR_MAP = {
    'red': COLOR_RED, 'orange': COLOR_ORANGE, 'blue': COLOR_BLUE,
    'green': COLOR_GREEN, 'gray': COLOR_GRAY, 'navy': COLOR_NAVY,
}

# ── 슬라이드 레이아웃 (TT_541 과 동일) ──────────────────────────
SLIDE_W = 13.33
SLIDE_H = 7.5
TITLE_H = 0.5
IMG_AREA_H = SLIDE_H - TITLE_H
DEFAULT_IMG_COL_W = 10.0
PDA_IMG_COL_W = 5.6
DESC_W_DEFAULT = SLIDE_W - DEFAULT_IMG_COL_W
FONT_NAME = '맑은 고딕'


def img_col_w_for(viewport):
    if viewport and viewport.get('height', 0) > viewport.get('width', 0):
        return PDA_IMG_COL_W
    return DEFAULT_IMG_COL_W


# ── 헬퍼 ────────────────────────────────────────────────────────
def png_size(p: Path):
    try:
        with Image.open(p) as im:
            return im.size
    except Exception:
        return None


def remove_all_slides(prs: Presentation):
    """템플릿 안의 기존 예제 슬라이드를 모두 제거 (마스터/레이아웃/테마는 보존)."""
    sldIdLst = prs.slides._sldIdLst
    pres_part = prs.part
    pres_rels = pres_part.rels
    targets = []
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        targets.append((sldId, rId))
    rels_internal = getattr(pres_rels, '_rels', None)
    for sldId, rId in targets:
        if rels_internal and rId in rels_internal:
            try:
                rels_internal.pop(rId)
            except Exception:
                pass
        sldIdLst.remove(sldId)


def set_text(tf, text: str, *, size=11, bold=False, color=COLOR_DARK,
             align=PP_ALIGN.LEFT, font=FONT_NAME):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_solid_rect(slide, x, y, w, h, fill=None, line=None, line_w=None,
                   no_fill=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.color.rgb = line if line else COLOR_WHITE
    if no_fill:
        sh.fill.background()
    elif fill is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        if line_w is not None:
            sh.line.width = Pt(line_w)
    sh.text_frame.text = ''
    sh.text_frame.margin_left = 0
    sh.text_frame.margin_right = 0
    sh.text_frame.margin_top = 0
    sh.text_frame.margin_bottom = 0
    return sh


def add_textbox(slide, x, y, w, h, text, *, size=11, bold=False,
                color=COLOR_DARK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                font=FONT_NAME):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    set_text(tf, text, size=size, bold=bold, color=color, align=align, font=font)
    return tb


def add_multipart_textbox(slide, x, y, w, h, parts, *, valign=MSO_ANCHOR.TOP,
                          font=FONT_NAME):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.clear()
    first = True
    for part in parts:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = part.get('align', PP_ALIGN.LEFT)
        if part.get('space_after'):
            p.space_after = Pt(part['space_after'])
        r = p.add_run()
        r.text = part.get('text', '')
        r.font.name = font
        r.font.size = Pt(part.get('size', 11))
        r.font.bold = bool(part.get('bold', False))
        r.font.color.rgb = part.get('color', COLOR_DARK)
    return tb


def add_line(slide, x1, y1, x2, y2, color, width_pt=1.0):
    cn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                    Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width_pt)
    return cn


# ── 좌표 계산: 픽셀 → 인치 ──────────────────────────────────────
class Geom:
    def __init__(self, px_w, px_h, img_col_w):
        scale = min(img_col_w / px_w, IMG_AREA_H / px_h)
        self.scale = scale
        self.disp_w = px_w * scale
        self.disp_h = px_h * scale
        self.img_l = (img_col_w - self.disp_w) / 2
        self.img_t = TITLE_H + (IMG_AREA_H - self.disp_h) / 2
        self.img_r = self.img_l + self.disp_w
        self.img_col_w = img_col_w

    def px_rect(self, px1, py1, px2, py2):
        return {
            'x': self.img_l + px1 * self.scale,
            'y': self.img_t + py1 * self.scale,
            'w': max(0.001, (px2 - px1) * self.scale),
            'h': max(0.001, (py2 - py1) * self.scale),
        }


# ── 영역 라벨링 (TT_541 과 동일) ────────────────────────────────
CIRCLED_NUMS = '①②③④⑤⑥⑦⑧⑨⑩⑪⑫'


def _region_number(label: str, idx: int) -> str:
    if label:
        if label[0] in CIRCLED_NUMS:
            return str(CIRCLED_NUMS.index(label[0]) + 1)
        m = re.match(r'^\s*(\d+)\s*[.)\s]', label)
        if m:
            return m.group(1)
    return str(idx + 1)


def add_region_labels(slide, regions, geom: Geom):
    BADGE_D = 0.42

    for idx, r in enumerate(regions):
        sl = geom.px_rect(*r['px'])
        color = r['color']
        num = _region_number(r.get('label', ''), idx)

        add_solid_rect(slide, sl['x'], sl['y'], sl['w'], sl['h'],
                       fill=None, line=color, line_w=2.0, no_fill=True)

        bx = sl['x'] - BADGE_D * 0.25
        by = sl['y'] - BADGE_D * 0.25
        bx = max(geom.img_l + 0.02, bx)
        by = max(geom.img_t + 0.02, by)

        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(bx), Inches(by), Inches(BADGE_D), Inches(BADGE_D),
        )
        oval.fill.solid()
        oval.fill.fore_color.rgb = color
        oval.line.color.rgb = COLOR_WHITE
        oval.line.width = Pt(1.75)
        tf = oval.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        n_len = len(num)
        text_size = 20 if n_len == 1 else (15 if n_len == 2 else 12)
        set_text(tf, num, size=text_size, bold=True, color=COLOR_WHITE,
                 align=PP_ALIGN.CENTER)


# ── 슬라이드 헤더 + 페이지 번호 ────────────────────────────────
def add_title_bar(slide, title_text):
    add_solid_rect(slide, 0, 0, SLIDE_W, TITLE_H,
                   fill=COLOR_TITLE_BG, line=COLOR_TITLE_BG, line_w=1.0)
    add_textbox(slide, 0.15, 0, SLIDE_W - 0.3, TITLE_H, title_text,
                size=16, bold=True, color=COLOR_WHITE,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)


def add_page_number(slide, idx, total):
    add_textbox(slide, SLIDE_W - 0.8, SLIDE_H - 0.28, 0.7, 0.22,
                f"{idx} / {total}", size=9, color=COLOR_PAGE_NUM,
                align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


# ── 인라인 [버튼명] 토큰 → 버튼 PNG 치환 (TT_541 동일) ─────────
BTN_TOKEN_RE = re.compile(r'\[([^\[\]]{1,16})\]')


def _norm_btn_name(s: str) -> str:
    return re.sub(r'\s+', '', s or '').strip()


_FONT_CACHE = {}
_FONT_CANDIDATES = [
    'C:/Windows/Fonts/malgun.ttf',
    'C:/Windows/Fonts/malgunsl.ttf',
    '/mnt/c/Windows/Fonts/malgun.ttf',
    '/mnt/c/Windows/Fonts/malgunsl.ttf',
    '/usr/share/fonts/truetype/nanum/NanumGothic.ttf',
    '/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc',
    '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
]


def _get_font(pt_size: float):
    key = round(pt_size, 1)
    if key in _FONT_CACHE:
        return _FONT_CACHE[key]
    px = max(6, int(round(pt_size * 96.0 / 72.0)))
    font = None
    for p in _FONT_CANDIDATES:
        try:
            if os.path.exists(p):
                font = ImageFont.truetype(p, px)
                break
        except Exception:
            continue
    _FONT_CACHE[key] = font
    return font


def _estimate_text_width(text: str, pt_size: float) -> float:
    if not text:
        return 0.0
    f = _get_font(pt_size)
    if f is not None:
        try:
            px_w = f.getlength(text)
            return float(px_w) / 96.0
        except Exception:
            pass
    w = 0.0
    for c in text:
        if c == ' ':
            w += pt_size * 0.005
        elif ord(c) > 0x7F and not c.isspace():
            w += pt_size * 0.015
        else:
            w += pt_size * 0.0075
    return w


def _add_seg_textbox(slide, x, y, w, h, text, *, size, bold, color):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    set_text(tf, text, size=size, bold=bold, color=color, align=PP_ALIGN.LEFT)
    return tb


def _line_segments(text: str, buttons: dict) -> list:
    if not buttons:
        return [{'kind': 'text', 'text': text}]
    btn_map = {_norm_btn_name(k): v for k, v in buttons.items()}
    segs = []
    last = 0
    for m in BTN_TOKEN_RE.finditer(text):
        norm = _norm_btn_name(m.group(1))
        if norm not in btn_map:
            continue
        if m.start() > last:
            segs.append({'kind': 'text', 'text': text[last:m.start()]})
        segs.append({'kind': 'img', 'path': btn_map[norm]})
        last = m.end()
    if last < len(text):
        segs.append({'kind': 'text', 'text': text[last:]})
    if not segs:
        segs = [{'kind': 'text', 'text': text}]
    return segs


def _render_line(slide, x0, y, max_w, segments, *,
                 text_size, text_color, text_bold,
                 line_h, img_h):
    cx = x0
    iy = y + (line_h - img_h) / 2
    GAP = 0.02
    for seg in segments:
        if cx - x0 >= max_w - 0.05:
            break
        if seg['kind'] == 'text':
            t = seg['text']
            if not t:
                continue
            tw = _estimate_text_width(t, text_size) + 0.03
            tw = min(tw, max_w - (cx - x0))
            if tw < 0.05:
                break
            _add_seg_textbox(slide, cx, y, tw, line_h, t,
                             size=text_size, bold=text_bold, color=text_color)
            cx += tw
        elif seg['kind'] == 'img':
            p = Path(seg['path'])
            sz = png_size(p) if p.exists() else None
            if not sz:
                continue
            iw_px, ih_px = sz
            iw = img_h * (iw_px / ih_px) if ih_px > 0 else img_h
            if cx - x0 + iw > max_w:
                iw = max_w - (cx - x0)
                if iw < 0.08:
                    break
            slide.shapes.add_picture(str(p), Inches(cx), Inches(iy),
                                     width=Inches(iw), height=Inches(img_h))
            cx += iw + GAP


def add_desc_panel(slide, desc_x, desc_w, desc_lines, buttons=None):
    buttons = buttons or {}
    add_solid_rect(slide, desc_x, TITLE_H, desc_w, IMG_AREA_H,
                   fill=COLOR_WHITE, line=COLOR_PANEL_LINE, line_w=1.0)

    x0 = desc_x + 0.12
    y = TITLE_H + 0.15
    max_w = desc_w - 0.24
    y_max = TITLE_H + IMG_AREA_H - 0.15

    HEADING_LINE_H = 0.30
    BODY_LINE_H = 0.26
    EMPTY_GAP = 0.08
    HEADING_AFTER = 0.04
    BODY_AFTER = 0.03
    HEADING_IMG_H = 0.22
    BODY_IMG_H = 0.22

    for entry in desc_lines:
        if y > y_max:
            break
        if entry == '' or entry is None:
            y += EMPTY_GAP
            continue
        if isinstance(entry, dict) and 'h' in entry:
            text = '■ ' + entry['h']
            color = entry.get('c', COLOR_NAVY)
            _render_line(slide, x0, y, max_w,
                         _line_segments(text, buttons),
                         text_size=13, text_color=color, text_bold=True,
                         line_h=HEADING_LINE_H, img_h=HEADING_IMG_H)
            y += HEADING_LINE_H + HEADING_AFTER
            continue
        if isinstance(entry, str) and entry.startswith('⚠'):
            _render_line(slide, x0, y, max_w,
                         _line_segments(entry, buttons),
                         text_size=11, text_color=COLOR_WARN, text_bold=False,
                         line_h=BODY_LINE_H, img_h=BODY_IMG_H)
            y += BODY_LINE_H + BODY_AFTER
            continue
        text = str(entry)
        _render_line(slide, x0, y, max_w,
                     _line_segments(text, buttons),
                     text_size=11, text_color=COLOR_DARK, text_bold=False,
                     line_h=BODY_LINE_H, img_h=BODY_IMG_H)
        y += BODY_LINE_H + BODY_AFTER


# ── ui.md 파싱 (선택) ──────────────────────────────────────────
def parse_ui_md(menu_code: str):
    p = REPO_ROOT / 'dist' / menu_code / 'ui.md'
    if not p.exists():
        return {'searchFields': [], 'gridColumns': [], 'rules': []}
    txt = p.read_text(encoding='utf-8', errors='replace')
    out = {'searchFields': [], 'gridColumns': [], 'rules': []}

    def section_bullets(name_re):
        m = re.search(name_re + r'[^\n]*\n([\s\S]*?)(?=\n##\s|$)', txt, re.M)
        if not m:
            return []
        bullets = re.findall(r'^[\-\*]\s+(.+)$', m.group(1), re.M)
        return [b.strip() for b in bullets if b.strip() and not re.match(r'^[-=]+$', b.strip())]

    out['searchFields'] = section_bullets(r'##\s*검색')[:8]
    out['gridColumns'] = section_bullets(r'##\s*(?:결과\s*)?그리드')[:12]
    out['rules'] = section_bullets(r'##\s*(?:공통\s*)?업무규칙')[:8]
    return out


# ── 슬라이드 빌더 ───────────────────────────────────────────────
def build_cover_slide(prs, customer):
    """표지 — 운영자 매뉴얼 (TT_541 과 다른 부분)"""
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    add_solid_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0xF4, 0xF6, 0xFA),
                   line=RGBColor(0xF4, 0xF6, 0xFA), line_w=0.5)
    add_solid_rect(s, 0, 2.4, SLIDE_W, 1.4, fill=COLOR_TITLE_BG,
                   line=COLOR_TITLE_BG, line_w=0.5)
    # ★ TT_541 과 다른 부분: 표지 제목 "운영자 매뉴얼"
    add_textbox(s, 0, 2.5, SLIDE_W, 0.7, '운영자 매뉴얼',
                size=40, bold=True, color=COLOR_WHITE,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 0, 3.15, SLIDE_W, 0.5, f'{customer} WMS',
                size=22, color=COLOR_SUBTITLE,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    today = datetime.date.today().isoformat()
    add_textbox(s, 0, 4.3, SLIDE_W, 0.4, f'작성일: {today}',
                size=13, color=RGBColor(0x56, 0x64, 0x76),
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return s


def build_toc_slide(prs, menus):
    """목차 — 카테고리별로 그룹핑"""
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    add_title_bar(s, '목 차')

    # 카테고리 → 메뉴 그룹핑
    groups = {}
    order = []
    for i, m in enumerate(menus):
        cat = m.get('category') or '운영'
        if cat not in groups:
            groups[cat] = []
            order.append(cat)
        groups[cat].append(m)

    parts = []
    idx = 1
    for cat in order:
        if len(order) > 1:
            parts.append({
                'text': f'[ {cat} ]',
                'size': 14, 'bold': True, 'color': COLOR_TITLE_BG, 'space_after': 6,
            })
        for m in groups[cat]:
            parts.append({
                'text': f"  {idx:02d}. {m['name']}  [{m['code'].upper()}]",
                'size': 14, 'color': COLOR_DARK, 'space_after': 4,
            })
            idx += 1
        parts.append({'text': '', 'size': 6, 'color': COLOR_DARK, 'space_after': 2})

    add_multipart_textbox(s, 0.6, 0.9, SLIDE_W - 1.2, SLIDE_H - 1.3, parts)
    return s


def build_category_section_slide(prs, category_name):
    """카테고리 섹션 표지 (시스템관리 / 마스터 / 권한 등)"""
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    add_solid_rect(s, 0, 0, SLIDE_W, SLIDE_H,
                   fill=RGBColor(0xE8, 0xEE, 0xF6),
                   line=RGBColor(0xE8, 0xEE, 0xF6), line_w=0.5)
    add_solid_rect(s, 0, 3.2, SLIDE_W, 1.0,
                   fill=COLOR_TITLE_BG, line=COLOR_TITLE_BG, line_w=0.5)
    add_textbox(s, 0, 3.25, SLIDE_W, 0.6, category_name,
                size=32, bold=True, color=COLOR_WHITE,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 0, 3.78, SLIDE_W, 0.4, '운영자 매뉴얼',
                size=15, color=COLOR_SUBTITLE,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return s


def build_menu_section_slide(prs, menu):
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    add_solid_rect(s, 0, 0, SLIDE_W, SLIDE_H,
                   fill=RGBColor(0xF4, 0xF6, 0xFA),
                   line=RGBColor(0xF4, 0xF6, 0xFA), line_w=0.5)
    add_solid_rect(s, 0, 3.0, SLIDE_W, 1.0,
                   fill=COLOR_TITLE_BG, line=COLOR_TITLE_BG, line_w=0.5)
    add_textbox(s, 0, 3.0, SLIDE_W, 0.6, menu['name'],
                size=36, bold=True, color=COLOR_WHITE,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 0, 3.55, SLIDE_W, 0.4, f"[{menu['code'].upper()}]",
                size=18, color=COLOR_SUBTITLE,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return s


def build_screen_slide(prs, *, title, image_path, geom, regions, desc_lines,
                       desc_x, desc_w, buttons=None):
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    add_title_bar(s, title)

    if image_path and image_path.exists():
        s.shapes.add_picture(
            str(image_path),
            Inches(geom.img_l), Inches(geom.img_t),
            width=Inches(geom.disp_w), height=Inches(geom.disp_h),
        )
    else:
        add_solid_rect(s, geom.img_l, geom.img_t, geom.disp_w, geom.disp_h,
                       fill=RGBColor(0xEE, 0xF2, 0xF7),
                       line=COLOR_PANEL_LINE, line_w=1.0)
        add_textbox(s, geom.img_l, geom.img_t, geom.disp_w, geom.disp_h,
                    '(이미지 없음)', size=16, color=RGBColor(0x90, 0xA4, 0xAE),
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    if regions:
        add_region_labels(s, regions, geom)

    add_desc_panel(s, desc_x, desc_w, desc_lines, buttons=buttons)
    return s


# ── 메뉴별 시나리오 → regions / desc 합성 (Vue 소스 기반) ───────
#
# 정책 변경 (사용자 요구사항):
#   - 01-main(메인 화면) 은 PPT 슬라이드에서 제외한다. 사용자가 보는 정보 가치가 낮고,
#     검색 실행 직후 결과 화면(02-search-result) 으로 충분.
#   - 우측 설명 패널은 일반 디폴트 문구가 아니라, Vue 소스에서 추출한 실제
#     검색필드 / 그리드 컬럼 / 버튼 / API 를 반영하여 작성한다.
#   - SearchSection 컴포넌트가 없거나 검색 필드가 비어 있는 메뉴는
#     "검색 조건" 라벨을 표시하지 않고 "기능 버튼 + 그리드" 형태로 표시한다.
SCENARIOS = [
    # 01-main 은 의도적으로 제외
    {'file': '02-search-result.png', 'kind': 'search', 'title': '메인 화면'},
    {'file': '03-register-popup.png', 'kind': 'register', 'title': '등록 팝업'},
    {'file': '04-row-selected.png', 'kind': 'rowSelect', 'title': '행 선택'},
    {'file': '05-edit-popup.png', 'kind': 'edit', 'title': '수정 팝업'},
]


def bbox_to_px(bb):
    if not bb or bb.get('width', 0) <= 0:
        return None
    return [bb['x'], bb['y'],
            bb['x'] + bb['width'], bb['y'] + bb['height']]


def _chunk_bullets(items, max_per_line=3, prefix='· '):
    """긴 컬럼 목록을 한 줄에 N개씩 묶어 가독성을 높인다."""
    if not items:
        return []
    lines = []
    cur = []
    for it in items:
        cur.append(it)
        if len(cur) >= max_per_line:
            lines.append(prefix + ', '.join(cur))
            cur = []
    if cur:
        lines.append(prefix + ', '.join(cur))
    return lines


def _filter_operator_buttons(buttons):
    """운영자 매뉴얼에 의미있는 버튼만 필터링 (PC매뉴얼/매뉴얼/로그아웃 등 제외)."""
    drop = {'PC매뉴얼', '매뉴얼', '로그아웃', '검색', '초기화', '닫기', '취소'}
    out = []
    for b in buttons:
        if not b or b in drop:
            continue
        if b in out:
            continue
        out.append(b)
    return out


def synth_regions_desc(kind, coords, vue):
    """Vue 소스에서 추출한 정보(vue)로 운영자 톤 설명을 생성한다.

    vue 는 parse_vue_source.parse_menu() 결과:
        has_search, search_fields, grid_columns, toolbar_buttons, apis 등.
    """
    regions, desc = [], []

    has_search = bool(vue.get('has_search') and vue.get('search_fields'))
    search_fields = vue.get('search_fields') or []
    grid_cols = vue.get('grid_columns') or []
    toolbar_buttons = _filter_operator_buttons(vue.get('toolbar_buttons') or [])

    if kind == 'search':
        # "메인 화면" 으로 표시되는 슬라이드 — 검색 후 결과 그리드 화면
        c = coords.get('search', {})
        gb = bbox_to_px(c.get('grid'))

        if has_search:
            # 검색 조건이 있는 메뉴
            if gb:
                regions.append({'px': gb, 'label': '1. 검색 결과 그리드', 'color': COLOR_GREEN})
            desc.append({'h': '1. 검색 영역', 'c': COLOR_RED})
            desc += _chunk_bullets(search_fields, max_per_line=3)
            desc.append('· 위 조건 입력 후 [검색] 버튼을 클릭하여 데이터를 조회합니다.')
            desc.append('')
            desc.append({'h': '2. 결과 그리드', 'c': COLOR_GREEN})
        else:
            # 검색 조건이 없는 메뉴 (행 직접 편집형)
            if gb:
                regions.append({'px': gb, 'label': '1. 그리드 영역', 'color': COLOR_GREEN})
            desc.append({'h': '1. 그리드 영역', 'c': COLOR_GREEN})
            desc.append('· 본 메뉴는 별도 검색 조건이 없으며, 진입 시 전체 목록이 자동 조회됩니다.')

        if grid_cols:
            desc.append('· 표시 컬럼:')
            desc += _chunk_bullets(grid_cols, max_per_line=3)

        if toolbar_buttons:
            desc.append('')
            desc.append({'h': '기능 버튼', 'c': COLOR_NAVY})
            desc += _chunk_bullets(toolbar_buttons, max_per_line=4, prefix='· ')

        desc.append('')
        desc.append({'h': '사용 방법', 'c': COLOR_NAVY})
        if vue.get('has_popup_edit'):
            desc.append('· 행 더블클릭 또는 [수정] 버튼 클릭 시 상세 편집 팝업이 열립니다.')
        else:
            desc.append('· 그리드 셀을 클릭해 인라인으로 값을 변경 후 [행저장] 버튼으로 저장합니다.')

    elif kind == 'register':
        c = coords.get('register', {})
        pp = bbox_to_px(c.get('popup'))
        if pp:
            regions.append({'px': pp, 'label': '1. 등록 팝업', 'color': COLOR_NAVY})
        desc.append({'h': '1. 등록 팝업', 'c': COLOR_NAVY})
        desc.append('· [추가] 버튼 클릭 시 신규 등록 팝업이 열립니다.')
        if grid_cols:
            desc.append('· 입력 항목은 그리드 컬럼과 동일합니다:')
            desc += _chunk_bullets(grid_cols[:8], max_per_line=3)
        desc.append('')
        desc.append({'h': '동작', 'c': COLOR_NAVY})
        desc.append('· 필수 항목(*)을 모두 채운 후 [저장] 버튼을 클릭하여 등록합니다.')
        desc.append('· [취소] 또는 [✕] 클릭 시 변경사항 없이 닫힙니다.')
        desc.append('')
        desc.append('⚠ 운영자 권한 변경은 시스템 전체에 영향을 주므로 신중하게 등록하세요.')
        desc.append('⚠ 본 매뉴얼은 팝업 화면을 보여주기 위해 열기만 한 상태입니다.')

    elif kind == 'rowSelect':
        c = coords.get('rowSelect', {})
        gb = bbox_to_px(c.get('grid'))
        if gb:
            regions.append({'px': gb, 'label': '1. 선택된 행', 'color': COLOR_BLUE})
        desc.append({'h': '1. 행 선택', 'c': COLOR_BLUE})
        desc.append('· 그리드에서 행을 클릭하면 해당 행이 강조되고 상세 정보가 조회됩니다.')
        if vue.get('has_popup_edit'):
            desc.append('· 선택된 행을 더블클릭하거나 [수정] 버튼을 클릭하면 상세 편집 팝업이 열립니다.')
        else:
            desc.append('· 선택된 행은 그리드 안에서 인라인으로 수정 가능합니다.')
        if toolbar_buttons:
            desc.append('')
            desc.append({'h': '실행 가능한 기능', 'c': COLOR_NAVY})
            desc += _chunk_bullets(toolbar_buttons, max_per_line=4)

    elif kind == 'edit':
        c = coords.get('edit', {})
        pp = bbox_to_px(c.get('popup'))
        if pp:
            regions.append({'px': pp, 'label': '1. 수정 팝업', 'color': COLOR_NAVY})
        desc.append({'h': '1. 수정 팝업', 'c': COLOR_NAVY})
        desc.append('· 결과 그리드에서 행 선택 후 [수정] 버튼 클릭 시 수정 팝업이 열립니다.')
        if grid_cols:
            desc.append('· 수정 가능 항목:')
            desc += _chunk_bullets(grid_cols[:8], max_per_line=3)
        desc.append('')
        desc.append({'h': '동작', 'c': COLOR_NAVY})
        desc.append('· 변경할 항목을 수정한 후 [저장] 버튼을 클릭하여 적용합니다.')
        desc.append('· [취소] 또는 [✕] 클릭 시 변경사항 없이 닫힙니다.')
        desc.append('')
        desc.append('⚠ 운영자 권한 변경은 시스템 전체에 영향을 주므로 신중하게 진행하세요.')
        desc.append('⚠ 본 매뉴얼은 팝업 화면을 보여주기 위해 열기만 한 상태입니다.')

    return regions, desc


# ── 메인 ────────────────────────────────────────────────────────
def main():
    prs = Presentation(str(TEMPLATE))
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    remove_all_slides(prs)

    # 표지 (운영자 매뉴얼)
    build_cover_slide(prs, CUSTOMER)
    # 목차 (카테고리별 그룹)
    build_toc_slide(prs, cfg['menus'])

    viewport = cfg.get('viewport', {'width': 1440, 'height': 900})
    img_col_w = img_col_w_for(viewport)
    desc_x = img_col_w
    desc_w = SLIDE_W - img_col_w

    # 카테고리별 분리 표지 (한 카테고리만 있으면 생략)
    distinct_categories = []
    for m in cfg['menus']:
        cat = m.get('category') or '운영'
        if cat not in distinct_categories:
            distinct_categories.append(cat)
    show_category_section = len(distinct_categories) >= 2

    fe_path = cfg.get('fePath')
    if not fe_path:
        print("[WARN] capture_config.json 에 fePath 가 없습니다. Vue 소스 파싱 생략.")

    last_category = None
    for menu in cfg['menus']:
        cat = menu.get('category') or '운영'
        if show_category_section and cat != last_category:
            build_category_section_slide(prs, cat)
            last_category = cat

        screen_dir = SCREENS_ROOT / menu['code']
        coords_path = screen_dir / 'coords.json'
        coords = (json.loads(coords_path.read_text(encoding='utf-8'))
                  if coords_path.exists() else {})

        # Vue 소스에서 실제 정보 추출 (검색필드/컬럼/버튼)
        if fe_path:
            try:
                vue = parse_vue_menu(menu['code'], fe_path)
            except Exception as e:
                print(f"[WARN] {menu['code']} Vue 파싱 실패: {e}")
                vue = {}
        else:
            vue = {}

        buttons_raw = coords.get('buttons') or {}
        buttons = {name: str(screen_dir / rel) for name, rel in buttons_raw.items()}

        # 메뉴 섹션 표지
        build_menu_section_slide(prs, menu)

        # 시나리오별 슬라이드 (01-main 은 SCENARIOS 에서 의도적으로 제외)
        for sc in SCENARIOS:
            img_path = screen_dir / sc['file']
            if not img_path.exists():
                continue
            size = png_size(img_path)
            px_w, px_h = size if size else (viewport['width'], viewport['height'])
            geom = Geom(px_w, px_h, img_col_w)

            regions, desc = synth_regions_desc(sc['kind'], coords, vue)
            title = f"{menu['name']} - {sc['title']}  [{menu['code'].upper()}]"
            build_screen_slide(prs, title=title, image_path=img_path,
                               geom=geom, regions=regions,
                               desc_lines=desc, desc_x=desc_x, desc_w=desc_w,
                               buttons=buttons)

    # 페이지 번호 (모든 슬라이드 마지막에 일괄)
    total = len(prs.slides)
    for i, sl in enumerate(prs.slides, start=1):
        add_page_number(sl, i, total)

    prs.save(str(OUT_FILE))

    print(f"[OK] 운영자매뉴얼 PPTX 생성 완료")
    print(f"     {OUT_FILE}")
    print(f"     총 슬라이드 {total}장 (메뉴 {len(cfg['menus'])}개, 카테고리 {len(distinct_categories)}개)")
    print(f"     템플릿     : {TEMPLATE}")


if __name__ == '__main__':
    main()
