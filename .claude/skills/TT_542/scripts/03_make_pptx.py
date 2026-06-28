#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
[TT_542] 4단계 — 사용자 매뉴얼 PPTX 생성 (템플릿 기반, python-pptx)

입력:
  deliverables/30-output/05 이행(TT)/tmp_542/capture_config.json
  deliverables/30-output/05 이행(TT)/tmp_542/screens/{메뉴코드}/*.png
  deliverables/30-output/05 이행(TT)/tmp_542/screens/{메뉴코드}/coords.json

템플릿:
  deliverables/10-templates/05 이행(TT)/TT.411-사용자매뉴얼(PDA).pptx
  - 13.33 × 7.5 인치 (16:9)
  - 제목바 #2D4B73 / 이미지영역 0~10in / 설명패널 10~13.33in / 페이지번호 우하단
  - 색상상수: 본 스크립트 상단의 COLOR_* 와 동일

출력:
  deliverables/30-output/05 이행(TT)/TT_542_사용자매뉴얼_PDA_{고객사명}.pptx

설계 원칙:
  - 라벨/테두리/배지/커넥터/설명패널은 모두 python-pptx 도형(`add_shape`)으로 PPT 안에
    직접 그린다. 이미지 합성하지 않으므로 PowerPoint 에서 자유롭게 편집 가능.
  - 이미지는 add_picture 의 width/height 직접 지정 (왜곡 방지).
  - 배지는 이미지 우측(IMG_R) ~ 설명 패널 사이 "배지 존" 에만 배치.
"""
from __future__ import annotations

import json
import os
import re
import sys
import datetime
from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image, ImageFont

# ── 경로 (스크립트 위치 기준 동적 계산) ────────────────────────
# 본 스크립트는 .claude/skills/TT_542/scripts/ 안에 있다.
# parents[0]=scripts  parents[1]=TT_542  parents[2]=skills  parents[3]=.claude  parents[4]=repo root
SCRIPT_PATH = Path(__file__).resolve()
REPO_ROOT = SCRIPT_PATH.parents[4]
TEMPLATE = REPO_ROOT / 'deliverables' / '10-templates' / '05 이행(TT)' / 'TT.411-사용자매뉴얼(PDA).pptx'
OUTPUT_DIR = REPO_ROOT / 'deliverables' / '30-output' / '05 이행(TT)'
TMP_DIR = OUTPUT_DIR / 'tmp_542'
SCREENS_ROOT = TMP_DIR / 'screens'
CFG_FILE = TMP_DIR / 'capture_config.json'

if not CFG_FILE.exists():
    sys.exit(f"[ERR] config 파일이 없습니다: {CFG_FILE}")
if not TEMPLATE.exists():
    sys.exit(f"[ERR] 템플릿 파일이 없습니다: {TEMPLATE}")

cfg = json.loads(CFG_FILE.read_text(encoding='utf-8'))


def sanitize_filename(s: str) -> str:
    return re.sub(r'[\\/:*?"<>|]', '_', str(s)).strip() or 'WMS'


CUSTOMER = sanitize_filename(cfg.get('customer') or 'WMS')
OUT_FILE = OUTPUT_DIR / f"TT_542_사용자매뉴얼_PDA_{CUSTOMER}.pptx"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ── 색상 상수 (템플릿 분석 결과와 동일) ─────────────────────────
COLOR_RED = RGBColor(0xDC, 0x1E, 0x1E)
COLOR_ORANGE = RGBColor(0xC8, 0x6E, 0x00)
COLOR_BLUE = RGBColor(0x1E, 0x64, 0xC8)
COLOR_GREEN = RGBColor(0x14, 0x8C, 0x3C)
COLOR_GRAY = RGBColor(0x6E, 0x6E, 0x6E)
COLOR_NAVY = RGBColor(0x1A, 0x3A, 0x5C)
COLOR_DARK = RGBColor(0x33, 0x33, 0x33)
COLOR_WARN = RGBColor(0xCC, 0x22, 0x22)
COLOR_TITLE_BG = RGBColor(0x2D, 0x4B, 0x73)
COLOR_PANEL_LINE = RGBColor(0xD0, 0xD5, 0xDD)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_PAGE_NUM = RGBColor(0x88, 0x88, 0x88)
COLOR_SUBTITLE = RGBColor(0xD6, 0xE0, 0xF0)

REGION_COLOR_MAP = {
    'red': COLOR_RED, 'orange': COLOR_ORANGE, 'blue': COLOR_BLUE,
    'green': COLOR_GREEN, 'gray': COLOR_GRAY, 'navy': COLOR_NAVY,
}

# ── 슬라이드 레이아웃 (템플릿과 동일) ───────────────────────────
SLIDE_W = 13.33
SLIDE_H = 7.5
TITLE_H = 0.5
IMG_AREA_H = SLIDE_H - TITLE_H  # 7.0
DEFAULT_IMG_COL_W = 10.0  # 템플릿 고정값 (가로형)
PDA_IMG_COL_W = 5.6       # PDA 세로 (대략 42%)
DESC_W_DEFAULT = SLIDE_W - DEFAULT_IMG_COL_W  # 3.33
FONT_NAME = '맑은 고딕'


def img_col_w_for(viewport):
    """뷰포트 비율로 이미지 컬럼 너비 결정 (PDA는 세로형 → 좁게)"""
    if viewport and viewport.get('height', 0) > viewport.get('width', 0):
        return PDA_IMG_COL_W
    return DEFAULT_IMG_COL_W


# ── 헬퍼 ────────────────────────────────────────────────────────
def png_size(path: Path):
    try:
        with Image.open(path) as im:
            return im.size  # (w, h)
    except Exception:
        return None


def remove_all_slides(prs: Presentation):
    """템플릿 안의 기존 예제 슬라이드를 모두 제거 (마스터/레이아웃/테마는 보존).

    python-pptx 의 _Relationships 는 직접 삭제를 지원하지 않으므로,
    내부 OXML 트리에서 직접 관계 엘리먼트와 sldId 를 제거한다.
    """
    sldIdLst = prs.slides._sldIdLst
    pres_part = prs.part
    pres_rels = pres_part.rels
    # 1) sldId 노드와 그 r:id 수집
    targets = []
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        targets.append((sldId, rId))
    # 2) 관계의 _rels 딕셔너리에서 직접 pop (내부 구현에 의존)
    rels_internal = getattr(pres_rels, '_rels', None)
    for sldId, rId in targets:
        if rels_internal and rId in rels_internal:
            try:
                rels_internal.pop(rId)
            except Exception:
                pass
        sldIdLst.remove(sldId)


def set_text(tf, text: str, *, size=11, bold=False, color=COLOR_DARK,
             align=PP_ALIGN.LEFT, font=FONT_NAME):
    """텍스트 프레임에 단일 단락을 세팅."""
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_solid_rect(slide, x, y, w, h, fill=None, line=None, line_w=None,
                   no_fill=False):
    """단순 사각형 도형 추가."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.color.rgb = line if line else COLOR_WHITE
    if no_fill:
        sh.fill.background()
    elif fill is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        if line_w is not None:
            sh.line.width = Pt(line_w)
    # 도형 자체 텍스트는 비움 (외부에서 별도 텍스트 박스로 처리)
    sh.text_frame.text = ''
    sh.text_frame.margin_left = 0
    sh.text_frame.margin_right = 0
    sh.text_frame.margin_top = 0
    sh.text_frame.margin_bottom = 0
    return sh


def add_textbox(slide, x, y, w, h, text, *, size=11, bold=False,
                color=COLOR_DARK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                font=FONT_NAME):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    set_text(tf, text, size=size, bold=bold, color=color, align=align, font=font)
    return tb


def add_multipart_textbox(slide, x, y, w, h, parts, *, valign=MSO_ANCHOR.TOP,
                          font=FONT_NAME):
    """parts: list of dict {text, size, bold, color, align}. 한 단락씩 추가."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.clear()
    first = True
    for part in parts:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = part.get('align', PP_ALIGN.LEFT)
        if part.get('space_after'):
            p.space_after = Pt(part['space_after'])
        r = p.add_run()
        r.text = part.get('text', '')
        r.font.name = font
        r.font.size = Pt(part.get('size', 11))
        r.font.bold = bool(part.get('bold', False))
        r.font.color.rgb = part.get('color', COLOR_DARK)
    return tb


def add_line(slide, x1, y1, x2, y2, color, width_pt=1.0):
    """직선(커넥터) 추가."""
    cn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                    Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width_pt)
    return cn


# ── 좌표 계산: 픽셀 → 인치 ──────────────────────────────────────
class Geom:
    def __init__(self, px_w, px_h, img_col_w):
        scale = min(img_col_w / px_w, IMG_AREA_H / px_h)
        self.scale = scale
        self.disp_w = px_w * scale
        self.disp_h = px_h * scale
        self.img_l = (img_col_w - self.disp_w) / 2
        self.img_t = TITLE_H + (IMG_AREA_H - self.disp_h) / 2
        self.img_r = self.img_l + self.disp_w
        self.img_col_w = img_col_w

    def px_rect(self, px1, py1, px2, py2):
        return {
            'x': self.img_l + px1 * self.scale,
            'y': self.img_t + py1 * self.scale,
            'w': max(0.001, (px2 - px1) * self.scale),
            'h': max(0.001, (py2 - py1) * self.scale),
        }


# ── 영역 라벨링 (테두리 + 좌상단 원형 번호 배지) ────────────────
CIRCLED_NUMS = '①②③④⑤⑥⑦⑧⑨⑩⑪⑫'


def _region_number(label: str, idx: int) -> str:
    """배지에 표시할 일반 숫자 문자열(1, 2, 3 ...).

    라벨 선두에 원형 숫자(①②...)나 `1.` `2.` 형식이 있으면 그 값을, 없으면 순번을 사용.
    """
    if label:
        if label[0] in CIRCLED_NUMS:
            return str(CIRCLED_NUMS.index(label[0]) + 1)
        m = re.match(r'^\s*(\d+)\s*[.)\s]', label)
        if m:
            return m.group(1)
    return str(idx + 1)


def add_region_labels(slide, regions, geom: Geom):
    """regions: list of {px:[x1,y1,x2,y2], label:str, color:RGBColor}

    각 영역에 색상 테두리 + 좌상단에 원형 번호 배지를 스크린샷 위에 직접 그린다.
    - 도형: MSO_SHAPE.OVAL (정원)
    - 채우기: 영역 색상 (region color)
    - 테두리: 흰색 1.75pt
    - 텍스트: 흰색 굵은 일반 숫자, 도형 크기에 맞춰 14~20pt 자동 스케일
    번호는 우측 설명 패널의 ■ 헤딩 번호와 1:1 매칭된다.
    """
    BADGE_D = 0.42  # 원형 배지 지름 (in)

    for idx, r in enumerate(regions):
        sl = geom.px_rect(*r['px'])
        color = r['color']
        num = _region_number(r.get('label', ''), idx)

        # 1) 영역 테두리 (투명 fill)
        add_solid_rect(slide, sl['x'], sl['y'], sl['w'], sl['h'],
                       fill=None, line=color, line_w=2.0, no_fill=True)

        # 2) 좌상단 원형 번호 배지 (이미지 위에 직접 배치)
        #    영역 좌상단 모서리에 1/4 정도 걸치도록 배치
        bx = sl['x'] - BADGE_D * 0.25
        by = sl['y'] - BADGE_D * 0.25
        # 이미지 영역을 벗어나지 않도록 보정
        bx = max(geom.img_l + 0.02, bx)
        by = max(geom.img_t + 0.02, by)

        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(bx), Inches(by), Inches(BADGE_D), Inches(BADGE_D),
        )
        oval.fill.solid()
        oval.fill.fore_color.rgb = color
        oval.line.color.rgb = COLOR_WHITE
        oval.line.width = Pt(1.75)
        tf = oval.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 자릿수에 따라 텍스트 사이즈 자동 스케일 (도형에 꽉 차게)
        n_len = len(num)
        text_size = 20 if n_len == 1 else (15 if n_len == 2 else 12)
        set_text(tf, num, size=text_size, bold=True, color=COLOR_WHITE,
                 align=PP_ALIGN.CENTER)


# ── 슬라이드 헤더(제목 바) + 페이지 번호 ────────────────────────
def add_title_bar(slide, title_text):
    add_solid_rect(slide, 0, 0, SLIDE_W, TITLE_H,
                   fill=COLOR_TITLE_BG, line=COLOR_TITLE_BG, line_w=1.0)
    add_textbox(slide, 0.15, 0, SLIDE_W - 0.3, TITLE_H, title_text,
                size=16, bold=True, color=COLOR_WHITE,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)


def add_page_number(slide, idx, total):
    add_textbox(slide, SLIDE_W - 0.8, SLIDE_H - 0.28, 0.7, 0.22,
                f"{idx} / {total}", size=9, color=COLOR_PAGE_NUM,
                align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


BTN_TOKEN_RE = re.compile(r'\[([^\[\]]{1,16})\]')


def _norm_btn_name(s: str) -> str:
    return re.sub(r'\s+', '', s or '').strip()


_FONT_CACHE = {}
_FONT_CANDIDATES = [
    r'C:\Windows\Fonts\malgun.ttf',
    r'C:\Windows\Fonts\malgunsl.ttf',
    '/mnt/c/Windows/Fonts/malgun.ttf',
    '/mnt/c/Windows/Fonts/malgunsl.ttf',
    '/usr/share/fonts/truetype/nanum/NanumGothic.ttf',
    '/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc',
    '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
]


def _get_font(pt_size: float):
    key = round(pt_size, 1)
    if key in _FONT_CACHE:
        return _FONT_CACHE[key]
    # PPT 인치 → PIL 픽셀: 96 DPI 기준 1pt = 96/72 px
    px = max(6, int(round(pt_size * 96.0 / 72.0)))
    font = None
    for p in _FONT_CANDIDATES:
        try:
            if os.path.exists(p):
                font = ImageFont.truetype(p, px)
                break
        except Exception:
            continue
    _FONT_CACHE[key] = font
    return font


def _estimate_text_width(text: str, pt_size: float) -> float:
    """텍스트 렌더 폭(in). PIL 실측 우선, 폰트 없으면 보수적 추정 fallback."""
    if not text:
        return 0.0
    f = _get_font(pt_size)
    if f is not None:
        try:
            px_w = f.getlength(text)
            return float(px_w) / 96.0
        except Exception:
            pass
    # Fallback: 보수적 추정
    w = 0.0
    for c in text:
        if c == ' ':
            w += pt_size * 0.005
        elif ord(c) > 0x7F and not c.isspace():
            w += pt_size * 0.015
        else:
            w += pt_size * 0.0075
    return w


def _add_seg_textbox(slide, x, y, w, h, text, *, size, bold, color):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    set_text(tf, text, size=size, bold=bold, color=color, align=PP_ALIGN.LEFT)
    return tb


def _line_segments(text: str, buttons: dict) -> list:
    """text 를 [{'kind':'text','text':...}, {'kind':'img','path':...}] 시퀀스로 분해.
    매칭되는 [버튼명] 토큰만 img 로 치환, 매칭 안 되면 그대로 텍스트로 유지.
    """
    if not buttons:
        return [{'kind': 'text', 'text': text}]
    btn_map = {_norm_btn_name(k): v for k, v in buttons.items()}
    segs = []
    last = 0
    for m in BTN_TOKEN_RE.finditer(text):
        norm = _norm_btn_name(m.group(1))
        if norm not in btn_map:
            continue
        if m.start() > last:
            segs.append({'kind': 'text', 'text': text[last:m.start()]})
        segs.append({'kind': 'img', 'path': btn_map[norm]})
        last = m.end()
    if last < len(text):
        segs.append({'kind': 'text', 'text': text[last:]})
    if not segs:
        segs = [{'kind': 'text', 'text': text}]
    return segs


def _render_line(slide, x0, y, max_w, segments, *,
                 text_size, text_color, text_bold,
                 line_h, img_h):
    """한 줄을 텍스트박스/이미지 가로 나열. 패널 너비 초과분은 잘림."""
    cx = x0
    iy = y + (line_h - img_h) / 2
    GAP = 0.02
    for seg in segments:
        if cx - x0 >= max_w - 0.05:
            break
        if seg['kind'] == 'text':
            t = seg['text']
            if not t:
                continue
            tw = _estimate_text_width(t, text_size) + 0.03
            tw = min(tw, max_w - (cx - x0))
            if tw < 0.05:
                break
            _add_seg_textbox(slide, cx, y, tw, line_h, t,
                             size=text_size, bold=text_bold, color=text_color)
            cx += tw
        elif seg['kind'] == 'img':
            p = Path(seg['path'])
            sz = png_size(p) if p.exists() else None
            if not sz:
                continue
            iw_px, ih_px = sz
            iw = img_h * (iw_px / ih_px) if ih_px > 0 else img_h
            if cx - x0 + iw > max_w:
                iw = max_w - (cx - x0)
                if iw < 0.08:
                    break
            slide.shapes.add_picture(str(p), Inches(cx), Inches(iy),
                                     width=Inches(iw), height=Inches(img_h))
            cx += iw + GAP


def add_desc_panel(slide, desc_x, desc_w, desc_lines, buttons=None):
    """desc_lines: list of dict / str
       - {'h': '헤딩', 'c': RGBColor} → 볼드 13pt 헤딩
       - '· 본문' → 11pt DARK (본문에 [버튼명] 토큰이 있으면 인라인 이미지로 치환)
       - '⚠ 경고' → 11pt WARN
       - '' → 빈 줄
    buttons: {버튼명: PNG 경로}. 매칭되는 토큰만 이미지로 치환.
    """
    buttons = buttons or {}

    # 패널 배경
    add_solid_rect(slide, desc_x, TITLE_H, desc_w, IMG_AREA_H,
                   fill=COLOR_WHITE, line=COLOR_PANEL_LINE, line_w=1.0)

    x0 = desc_x + 0.12
    y = TITLE_H + 0.15
    max_w = desc_w - 0.24
    y_max = TITLE_H + IMG_AREA_H - 0.15

    HEADING_LINE_H = 0.30
    BODY_LINE_H = 0.26
    EMPTY_GAP = 0.08
    HEADING_AFTER = 0.04
    BODY_AFTER = 0.03
    HEADING_IMG_H = 0.22
    BODY_IMG_H = 0.22

    for entry in desc_lines:
        if y > y_max:
            break
        if entry == '' or entry is None:
            y += EMPTY_GAP
            continue
        if isinstance(entry, dict) and 'h' in entry:
            text = '■ ' + entry['h']
            color = entry.get('c', COLOR_NAVY)
            _render_line(slide, x0, y, max_w,
                         _line_segments(text, buttons),
                         text_size=13, text_color=color, text_bold=True,
                         line_h=HEADING_LINE_H, img_h=HEADING_IMG_H)
            y += HEADING_LINE_H + HEADING_AFTER
            continue
        if isinstance(entry, str) and entry.startswith('⚠'):
            _render_line(slide, x0, y, max_w,
                         _line_segments(entry, buttons),
                         text_size=11, text_color=COLOR_WARN, text_bold=False,
                         line_h=BODY_LINE_H, img_h=BODY_IMG_H)
            y += BODY_LINE_H + BODY_AFTER
            continue
        text = str(entry)
        _render_line(slide, x0, y, max_w,
                     _line_segments(text, buttons),
                     text_size=11, text_color=COLOR_DARK, text_bold=False,
                     line_h=BODY_LINE_H, img_h=BODY_IMG_H)
        y += BODY_LINE_H + BODY_AFTER


# ── ui.md 파싱 (선택) ──────────────────────────────────────────
def parse_ui_md(menu_code: str):
    p = REPO_ROOT / 'dist' / menu_code / 'ui.md'
    if not p.exists():
        return {'searchFields': [], 'gridColumns': [], 'rules': []}
    txt = p.read_text(encoding='utf-8', errors='replace')
    out = {'searchFields': [], 'gridColumns': [], 'rules': []}

    def section_bullets(name_re):
        m = re.search(name_re + r'[^\n]*\n([\s\S]*?)(?=\n##\s|$)', txt, re.M)
        if not m:
            return []
        bullets = re.findall(r'^[\-\*]\s+(.+)$', m.group(1), re.M)
        return [b.strip() for b in bullets if b.strip() and not re.match(r'^[-=]+$', b.strip())]

    out['searchFields'] = section_bullets(r'##\s*검색')[:8]
    out['gridColumns'] = section_bullets(r'##\s*(?:결과\s*)?그리드')[:12]
    out['rules'] = section_bullets(r'##\s*(?:공통\s*)?업무규칙')[:8]
    return out


# ── 슬라이드 빌더 ───────────────────────────────────────────────
def build_cover_slide(prs, customer):
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    # 배경 영역 (메인)
    add_solid_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0xF4, 0xF6, 0xFA),
                   line=RGBColor(0xF4, 0xF6, 0xFA), line_w=0.5)
    # 중앙 띠
    add_solid_rect(s, 0, 2.4, SLIDE_W, 1.4, fill=COLOR_TITLE_BG,
                   line=COLOR_TITLE_BG, line_w=0.5)
    add_textbox(s, 0, 2.5, SLIDE_W, 0.7, '사용자 매뉴얼 (PDA)',
                size=40, bold=True, color=COLOR_WHITE,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 0, 3.15, SLIDE_W, 0.5, f'{customer} WMS',
                size=22, color=COLOR_SUBTITLE,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    today = datetime.date.today().isoformat()
    add_textbox(s, 0, 4.3, SLIDE_W, 0.4, f'작성일: {today}',
                size=13, color=RGBColor(0x56, 0x64, 0x76),
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return s


def build_toc_slide(prs, menus):
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    add_title_bar(s, '목 차')
    lines = [f"{(i + 1):02d}. {m['name']}  [{m['code'].upper()}]"
             for i, m in enumerate(menus)]
    parts = [{'text': line, 'size': 16, 'color': COLOR_DARK, 'space_after': 8}
             for line in lines]
    add_multipart_textbox(s, 0.6, 0.9, SLIDE_W - 1.2, SLIDE_H - 1.3, parts)
    return s


def build_menu_section_slide(prs, menu):
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    add_solid_rect(s, 0, 0, SLIDE_W, SLIDE_H,
                   fill=RGBColor(0xF4, 0xF6, 0xFA),
                   line=RGBColor(0xF4, 0xF6, 0xFA), line_w=0.5)
    add_solid_rect(s, 0, 3.0, SLIDE_W, 1.0,
                   fill=COLOR_TITLE_BG, line=COLOR_TITLE_BG, line_w=0.5)
    add_textbox(s, 0, 3.0, SLIDE_W, 0.6, menu['name'],
                size=36, bold=True, color=COLOR_WHITE,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 0, 3.55, SLIDE_W, 0.4, f"[{menu['code'].upper()}]",
                size=18, color=COLOR_SUBTITLE,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return s


def build_screen_slide(prs, *, title, image_path, geom, regions, desc_lines,
                       desc_x, desc_w, buttons=None):
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    add_title_bar(s, title)

    if image_path and image_path.exists():
        s.shapes.add_picture(
            str(image_path),
            Inches(geom.img_l), Inches(geom.img_t),
            width=Inches(geom.disp_w), height=Inches(geom.disp_h),
        )
    else:
        # 이미지 없음 placeholder
        add_solid_rect(s, geom.img_l, geom.img_t, geom.disp_w, geom.disp_h,
                       fill=RGBColor(0xEE, 0xF2, 0xF7),
                       line=COLOR_PANEL_LINE, line_w=1.0)
        add_textbox(s, geom.img_l, geom.img_t, geom.disp_w, geom.disp_h,
                    '(이미지 없음)', size=16, color=RGBColor(0x90, 0xA4, 0xAE),
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    if regions:
        add_region_labels(s, regions, geom)

    add_desc_panel(s, desc_x, desc_w, desc_lines, buttons=buttons)
    return s


# ── 메뉴별 시나리오 → regions / desc 합성 ──────────────────────
SCENARIOS = [
    {'file': '01-main.png', 'kind': 'main', 'title': '메인 화면'},
    {'file': '02-search-result.png', 'kind': 'search', 'title': '검색 결과'},
    {'file': '03-register-popup.png', 'kind': 'register', 'title': '등록 팝업'},
    {'file': '04-row-selected.png', 'kind': 'rowSelect', 'title': '행 선택'},
    {'file': '05-edit-popup.png', 'kind': 'edit', 'title': '수정 팝업'},
]


def bbox_to_px(bb):
    if not bb or bb.get('width', 0) <= 0:
        return None
    return [bb['x'], bb['y'],
            bb['x'] + bb['width'], bb['y'] + bb['height']]


def synth_regions_desc(kind, coords, ui):
    regions, desc = [], []
    if kind == 'main':
        c = coords.get('main', {})
        sb = bbox_to_px(c.get('search'))
        gb = bbox_to_px(c.get('grid'))
        if sb:
            regions.append({'px': sb, 'label': '1. 검색 조건', 'color': COLOR_RED})
            desc.append({'h': '1. 검색 조건', 'c': COLOR_RED})
            if ui['searchFields']:
                desc += ['· ' + f for f in ui['searchFields']]
            else:
                desc.append('· 화면 상단 조건을 입력 후 [검색] 클릭')
            desc.append('')
        if gb:
            regions.append({'px': gb, 'label': '2. 결과 그리드', 'color': COLOR_GRAY})
            desc.append({'h': '2. 결과 그리드', 'c': COLOR_GRAY})
            desc.append('· 검색 실행 전 빈 상태로 표시됩니다.')
            desc.append('')
        if ui['rules']:
            desc.append({'h': '업무 규칙', 'c': COLOR_NAVY})
            desc += ['· ' + r for r in ui['rules']]
    elif kind == 'search':
        c = coords.get('search', {})
        gb = bbox_to_px(c.get('grid'))
        if gb:
            regions.append({'px': gb, 'label': '1. 검색 결과 그리드', 'color': COLOR_GREEN})
            desc.append({'h': '1. 검색 결과 그리드', 'c': COLOR_GREEN})
        else:
            desc.append({'h': '검색 결과', 'c': COLOR_GREEN})
        desc.append('· 입력한 조건에 해당하는 데이터가 표시됩니다.')
        if ui['gridColumns']:
            desc.append('')
            desc.append({'h': '주요 컬럼', 'c': COLOR_NAVY})
            desc += ['· ' + col for col in ui['gridColumns']]
        desc.append('')
        desc.append({'h': '사용 방법', 'c': COLOR_NAVY})
        desc.append('· 결과 행을 클릭하면 상세 정보가 표시됩니다.')
        desc.append('· 페이지 하단에서 페이지 이동 / 페이지 크기를 변경할 수 있습니다.')
    elif kind in ('register', 'edit'):
        key = 'register' if kind == 'register' else 'edit'
        c = coords.get(key, {})
        pp = bbox_to_px(c.get('popup'))
        label = '등록 팝업' if kind == 'register' else '수정 팝업'
        numbered = '1. ' + label
        if pp:
            regions.append({'px': pp, 'label': numbered, 'color': COLOR_NAVY})
            desc.append({'h': numbered, 'c': COLOR_NAVY})
        else:
            desc.append({'h': label, 'c': COLOR_NAVY})
        if kind == 'register':
            desc.append('· [추가] 버튼 클릭 시 표시됩니다.')
        else:
            desc.append('· 결과 그리드에서 행 선택 후 [수정] 클릭 시 표시됩니다.')
        desc.append('· 필수 항목(*)을 모두 입력한 후 [저장] 버튼을 클릭합니다.')
        desc.append('· [취소] 또는 [✕] 클릭 시 변경사항 없이 닫힙니다.')
        desc.append('')
        desc.append('⚠ 본 매뉴얼은 팝업 화면을 보여주기 위해 열기만 한 상태입니다.')
    elif kind == 'rowSelect':
        c = coords.get('rowSelect', {})
        gb = bbox_to_px(c.get('grid'))
        if gb:
            regions.append({'px': gb, 'label': '1. 선택된 행', 'color': COLOR_BLUE})
            desc.append({'h': '1. 선택된 행', 'c': COLOR_BLUE})
        else:
            desc.append({'h': '행 선택', 'c': COLOR_BLUE})
        desc.append('· 그리드 행 클릭 시 해당 행이 강조 표시됩니다.')
        desc.append('· 행 선택 후 [수정] / [삭제] / [복사] 등의 기능을 실행할 수 있습니다.')
    return regions, desc


# ── 메인 ────────────────────────────────────────────────────────
def main():
    prs = Presentation(str(TEMPLATE))
    # 슬라이드 크기 보정 (안전상)
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # 템플릿의 예제 슬라이드는 모두 제거 (마스터/레이아웃/테마는 유지)
    remove_all_slides(prs)

    # 표지
    build_cover_slide(prs, CUSTOMER)
    # 목차
    build_toc_slide(prs, cfg['menus'])

    viewport = cfg.get('viewport', {'width': 1440, 'height': 900})
    img_col_w = img_col_w_for(viewport)
    desc_x = img_col_w
    desc_w = SLIDE_W - img_col_w

    # 메뉴별 슬라이드
    for menu in cfg['menus']:
        screen_dir = SCREENS_ROOT / menu['code']
        coords_path = screen_dir / 'coords.json'
        coords = (json.loads(coords_path.read_text(encoding='utf-8'))
                  if coords_path.exists() else {})
        ui = parse_ui_md(menu['code'])
        # 버튼 PNG 경로 (coords.buttons 의 상대경로 → 절대경로 매핑)
        buttons_raw = coords.get('buttons') or {}
        buttons = {name: str(screen_dir / rel) for name, rel in buttons_raw.items()}

        # 메뉴 섹션 표지
        build_menu_section_slide(prs, menu)

        # 시나리오별 슬라이드
        for sc in SCENARIOS:
            img_path = screen_dir / sc['file']
            if not img_path.exists():
                continue
            size = png_size(img_path)
            px_w, px_h = size if size else (viewport['width'], viewport['height'])
            geom = Geom(px_w, px_h, img_col_w)

            regions, desc = synth_regions_desc(sc['kind'], coords, ui)
            title = f"{menu['name']} - {sc['title']}  [{menu['code'].upper()}]"
            build_screen_slide(prs, title=title, image_path=img_path,
                               geom=geom, regions=regions,
                               desc_lines=desc, desc_x=desc_x, desc_w=desc_w,
                               buttons=buttons)

    # 페이지 번호 (모든 슬라이드 마지막에 일괄)
    total = len(prs.slides)
    for i, sl in enumerate(prs.slides, start=1):
        add_page_number(sl, i, total)

    prs.save(str(OUT_FILE))

    print(f"[OK] PPTX 생성 완료")
    print(f"     {OUT_FILE}")
    print(f"     총 슬라이드 {total}장 (메뉴 {len(cfg['menus'])}개)")
    print(f"     템플릿     : {TEMPLATE}")


if __name__ == '__main__':
    main()
