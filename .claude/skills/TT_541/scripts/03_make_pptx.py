#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
[TT_541] 4단계 — 사용자 매뉴얼 PPTX 생성 (템플릿 기반, python-pptx)

입력:
  output/05 이행(TT)/tmp/capture_config.json
  output/05 이행(TT)/tmp/screens/{메뉴코드}/*.png
  output/05 이행(TT)/tmp/screens/{메뉴코드}/coords.json

템플릿:
  template/05 이행(TT)/사용자_매뉴얼_템플릿.pptx
  - 13.33 × 7.5 인치 (16:9)
  - 제목바 #2D4B73 / 이미지영역 0~10in / 설명패널 10~13.33in / 페이지번호 우하단
  - 색상상수: 본 스크립트 상단의 COLOR_* 와 동일

출력:
  output/05 이행(TT)/TT_541_사용자매뉴얼_{고객사명}.pptx

설계 원칙:
  - 라벨/테두리/배지/커넥터/설명패널은 모두 python-pptx 도형(`add_shape`)으로 PPT 안에
    직접 그린다. 이미지 합성하지 않으므로 PowerPoint 에서 자유롭게 편집 가능.
  - 이미지는 add_picture 의 width/height 직접 지정 (왜곡 방지).
  - 배지는 이미지 우측(IMG_R) ~ 설명 패널 사이 "배지 존" 에만 배치.
"""
from __future__ import annotations

import json
import os
import re
import sys
import datetime
from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

# ── 경로 ────────────────────────────────────────────────────────
REPO_ROOT = Path('/mnt/c/zinide/workspace/cloud-wms-doc')
TEMPLATE = REPO_ROOT / 'template' / '05 이행(TT)' / '사용자_매뉴얼_템플릿.pptx'
OUTPUT_DIR = REPO_ROOT / 'output' / '05 이행(TT)'
TMP_DIR = OUTPUT_DIR / 'tmp'
SCREENS_ROOT = TMP_DIR / 'screens'
CFG_FILE = TMP_DIR / 'capture_config.json'

if not CFG_FILE.exists():
    sys.exit(f"[ERR] config 파일이 없습니다: {CFG_FILE}")
if not TEMPLATE.exists():
    sys.exit(f"[ERR] 템플릿 파일이 없습니다: {TEMPLATE}")

cfg = json.loads(CFG_FILE.read_text(encoding='utf-8'))


def sanitize_filename(s: str) -> str:
    return re.sub(r'[\\/:*?"<>|]', '_', str(s)).strip() or 'WMS'


CUSTOMER = sanitize_filename(cfg.get('customer') or 'WMS')
OUT_FILE = OUTPUT_DIR / f"TT_541_사용자매뉴얼_{CUSTOMER}.pptx"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ── 색상 상수 (템플릿 분석 결과와 동일) ─────────────────────────
COLOR_RED = RGBColor(0xDC, 0x1E, 0x1E)
COLOR_ORANGE = RGBColor(0xC8, 0x6E, 0x00)
COLOR_BLUE = RGBColor(0x1E, 0x64, 0xC8)
COLOR_GREEN = RGBColor(0x14, 0x8C, 0x3C)
COLOR_GRAY = RGBColor(0x6E, 0x6E, 0x6E)
COLOR_NAVY = RGBColor(0x1A, 0x3A, 0x5C)
COLOR_DARK = RGBColor(0x33, 0x33, 0x33)
COLOR_WARN = RGBColor(0xCC, 0x22, 0x22)
COLOR_TITLE_BG = RGBColor(0x2D, 0x4B, 0x73)
COLOR_PANEL_LINE = RGBColor(0xD0, 0xD5, 0xDD)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_PAGE_NUM = RGBColor(0x88, 0x88, 0x88)
COLOR_SUBTITLE = RGBColor(0xD6, 0xE0, 0xF0)

REGION_COLOR_MAP = {
    'red': COLOR_RED, 'orange': COLOR_ORANGE, 'blue': COLOR_BLUE,
    'green': COLOR_GREEN, 'gray': COLOR_GRAY, 'navy': COLOR_NAVY,
}

# ── 슬라이드 레이아웃 (템플릿과 동일) ───────────────────────────
SLIDE_W = 13.33
SLIDE_H = 7.5
TITLE_H = 0.5
IMG_AREA_H = SLIDE_H - TITLE_H  # 7.0
DEFAULT_IMG_COL_W = 10.0  # 템플릿 고정값 (가로형)
PDA_IMG_COL_W = 5.6       # PDA 세로 (대략 42%)
DESC_W_DEFAULT = SLIDE_W - DEFAULT_IMG_COL_W  # 3.33
FONT_NAME = '맑은 고딕'


def img_col_w_for(viewport):
    """뷰포트 비율로 이미지 컬럼 너비 결정 (PDA는 세로형 → 좁게)"""
    if viewport and viewport.get('height', 0) > viewport.get('width', 0):
        return PDA_IMG_COL_W
    return DEFAULT_IMG_COL_W


# ── 헬퍼 ────────────────────────────────────────────────────────
def png_size(path: Path):
    try:
        with Image.open(path) as im:
            return im.size  # (w, h)
    except Exception:
        return None


def remove_all_slides(prs: Presentation):
    """템플릿 안의 기존 예제 슬라이드를 모두 제거 (마스터/레이아웃/테마는 보존).

    python-pptx 의 _Relationships 는 직접 삭제를 지원하지 않으므로,
    내부 OXML 트리에서 직접 관계 엘리먼트와 sldId 를 제거한다.
    """
    sldIdLst = prs.slides._sldIdLst
    pres_part = prs.part
    pres_rels = pres_part.rels
    # 1) sldId 노드와 그 r:id 수집
    targets = []
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        targets.append((sldId, rId))
    # 2) 관계의 _rels 딕셔너리에서 직접 pop (내부 구현에 의존)
    rels_internal = getattr(pres_rels, '_rels', None)
    for sldId, rId in targets:
        if rels_internal and rId in rels_internal:
            try:
                rels_internal.pop(rId)
            except Exception:
                pass
        sldIdLst.remove(sldId)


def set_text(tf, text: str, *, size=11, bold=False, color=COLOR_DARK,
             align=PP_ALIGN.LEFT, font=FONT_NAME):
    """텍스트 프레임에 단일 단락을 세팅."""
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_solid_rect(slide, x, y, w, h, fill=None, line=None, line_w=None,
                   no_fill=False):
    """단순 사각형 도형 추가."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.color.rgb = line if line else COLOR_WHITE
    if no_fill:
        sh.fill.background()
    elif fill is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        if line_w is not None:
            sh.line.width = Pt(line_w)
    # 도형 자체 텍스트는 비움 (외부에서 별도 텍스트 박스로 처리)
    sh.text_frame.text = ''
    sh.text_frame.margin_left = 0
    sh.text_frame.margin_right = 0
    sh.text_frame.margin_top = 0
    sh.text_frame.margin_bottom = 0
    return sh


def add_textbox(slide, x, y, w, h, text, *, size=11, bold=False,
                color=COLOR_DARK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                font=FONT_NAME):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    set_text(tf, text, size=size, bold=bold, color=color, align=align, font=font)
    return tb


def add_multipart_textbox(slide, x, y, w, h, parts, *, valign=MSO_ANCHOR.TOP,
                          font=FONT_NAME):
    """parts: list of dict {text, size, bold, color, align}. 한 단락씩 추가."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.clear()
    first = True
    for part in parts:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = part.get('align', PP_ALIGN.LEFT)
        if part.get('space_after'):
            p.space_after = Pt(part['space_after'])
        r = p.add_run()
        r.text = part.get('text', '')
        r.font.name = font
        r.font.size = Pt(part.get('size', 11))
        r.font.bold = bool(part.get('bold', False))
        r.font.color.rgb = part.get('color', COLOR_DARK)
    return tb


def add_line(slide, x1, y1, x2, y2, color, width_pt=1.0):
    """직선(커넥터) 추가."""
    cn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                    Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width_pt)
    return cn


# ── 좌표 계산: 픽셀 → 인치 ──────────────────────────────────────
class Geom:
    def __init__(self, px_w, px_h, img_col_w):
        scale = min(img_col_w / px_w, IMG_AREA_H / px_h)
        self.scale = scale
        self.disp_w = px_w * scale
        self.disp_h = px_h * scale
        self.img_l = (img_col_w - self.disp_w) / 2
        self.img_t = TITLE_H + (IMG_AREA_H - self.disp_h) / 2
        self.img_r = self.img_l + self.disp_w
        self.img_col_w = img_col_w

    def px_rect(self, px1, py1, px2, py2):
        return {
            'x': self.img_l + px1 * self.scale,
            'y': self.img_t + py1 * self.scale,
            'w': max(0.001, (px2 - px1) * self.scale),
            'h': max(0.001, (py2 - py1) * self.scale),
        }


# ── 영역 라벨링 (테두리 + 배지 + 커넥터) ────────────────────────
def add_region_labels(slide, regions, geom: Geom):
    """regions: list of {px:[x1,y1,x2,y2], label:str, color:RGBColor}"""
    BADGE_H = 0.22
    BADGE_X = geom.img_r + 0.04
    BADGE_W = max(0.4, geom.img_col_w - geom.img_r - 0.08)

    badges = []
    for r in regions:
        sl = geom.px_rect(*r['px'])
        mid_y = sl['y'] + sl['h'] / 2
        badges.append({
            **r, 'sl': sl, 'mid_y': mid_y,
            'badge_y': mid_y - BADGE_H / 2,
        })
    badges.sort(key=lambda b: b['badge_y'])
    for i in range(1, len(badges)):
        prev = badges[i - 1]
        cur = badges[i]
        if cur['badge_y'] < prev['badge_y'] + BADGE_H + 0.04:
            cur['badge_y'] = prev['badge_y'] + BADGE_H + 0.04

    for b in badges:
        sl = b['sl']
        color = b['color']

        # 1) 영역 테두리 (투명 fill)
        add_solid_rect(slide, sl['x'], sl['y'], sl['w'], sl['h'],
                       fill=None, line=color, line_w=2.0, no_fill=True)

        # 2) 커넥터 (영역 우측 중앙 → 배지 좌측 중앙)
        badge_mid_y = b['badge_y'] + BADGE_H / 2
        add_line(slide, geom.img_r, b['mid_y'], BADGE_X, badge_mid_y,
                 color=color, width_pt=1.0)

        # 3) 배지 사각형
        add_solid_rect(slide, BADGE_X, b['badge_y'], BADGE_W, BADGE_H,
                       fill=color, line=color, line_w=0.5)

        # 4) 배지 텍스트
        add_textbox(slide, BADGE_X, b['badge_y'], BADGE_W, BADGE_H,
                    b['label'], size=8.5, bold=True, color=COLOR_WHITE,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


# ── 슬라이드 헤더(제목 바) + 페이지 번호 ────────────────────────
def add_title_bar(slide, title_text):
    add_solid_rect(slide, 0, 0, SLIDE_W, TITLE_H,
                   fill=COLOR_TITLE_BG, line=COLOR_TITLE_BG, line_w=1.0)
    add_textbox(slide, 0.15, 0, SLIDE_W - 0.3, TITLE_H, title_text,
                size=16, bold=True, color=COLOR_WHITE,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)


def add_page_number(slide, idx, total):
    add_textbox(slide, SLIDE_W - 0.8, SLIDE_H - 0.28, 0.7, 0.22,
                f"{idx} / {total}", size=9, color=COLOR_PAGE_NUM,
                align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


def add_desc_panel(slide, desc_x, desc_w, desc_lines):
    """desc_lines: list of dict / str
       - {'h': '헤딩', 'c': RGBColor} → 볼드 13pt 헤딩
       - '· 본문' → 11pt DARK
       - '⚠ 경고' → 11pt WARN
       - '' → 빈 줄
    """
    # 패널 배경
    add_solid_rect(slide, desc_x, TITLE_H, desc_w, IMG_AREA_H,
                   fill=COLOR_WHITE, line=COLOR_PANEL_LINE, line_w=1.0)
    parts = []
    for entry in desc_lines:
        if entry == '' or entry is None:
            parts.append({'text': ' ', 'size': 5, 'space_after': 2})
            continue
        if isinstance(entry, dict) and 'h' in entry:
            parts.append({
                'text': '■ ' + entry['h'],
                'size': 13, 'bold': True,
                'color': entry.get('c', COLOR_NAVY),
                'space_after': 4,
            })
            continue
        if isinstance(entry, str) and entry.startswith('⚠'):
            parts.append({
                'text': entry, 'size': 11, 'color': COLOR_WARN,
                'space_after': 3,
            })
            continue
        parts.append({
            'text': str(entry), 'size': 11, 'color': COLOR_DARK,
            'space_after': 3,
        })
    add_multipart_textbox(slide, desc_x + 0.12, TITLE_H + 0.15,
                          desc_w - 0.24, IMG_AREA_H - 0.3, parts,
                          valign=MSO_ANCHOR.TOP)


# ── ui.md 파싱 (선택) ──────────────────────────────────────────
def parse_ui_md(menu_code: str):
    p = REPO_ROOT / 'dist' / menu_code / 'ui.md'
    if not p.exists():
        return {'searchFields': [], 'gridColumns': [], 'rules': []}
    txt = p.read_text(encoding='utf-8', errors='replace')
    out = {'searchFields': [], 'gridColumns': [], 'rules': []}

    def section_bullets(name_re):
        m = re.search(name_re + r'[^\n]*\n([\s\S]*?)(?=\n##\s|$)', txt, re.M)
        if not m:
            return []
        bullets = re.findall(r'^[\-\*]\s+(.+)$', m.group(1), re.M)
        return [b.strip() for b in bullets if b.strip() and not re.match(r'^[-=]+$', b.strip())]

    out['searchFields'] = section_bullets(r'##\s*검색')[:8]
    out['gridColumns'] = section_bullets(r'##\s*(?:결과\s*)?그리드')[:12]
    out['rules'] = section_bullets(r'##\s*(?:공통\s*)?업무규칙')[:8]
    return out


# ── 슬라이드 빌더 ───────────────────────────────────────────────
def build_cover_slide(prs, customer):
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    # 배경 영역 (메인)
    add_solid_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0xF4, 0xF6, 0xFA),
                   line=RGBColor(0xF4, 0xF6, 0xFA), line_w=0.5)
    # 중앙 띠
    add_solid_rect(s, 0, 2.4, SLIDE_W, 1.4, fill=COLOR_TITLE_BG,
                   line=COLOR_TITLE_BG, line_w=0.5)
    add_textbox(s, 0, 2.5, SLIDE_W, 0.7, '사용자 매뉴얼',
                size=40, bold=True, color=COLOR_WHITE,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 0, 3.15, SLIDE_W, 0.5, f'{customer} WMS',
                size=22, color=COLOR_SUBTITLE,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    today = datetime.date.today().isoformat()
    add_textbox(s, 0, 4.3, SLIDE_W, 0.4, f'작성일: {today}',
                size=13, color=RGBColor(0x56, 0x64, 0x76),
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return s


def build_toc_slide(prs, menus):
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    add_title_bar(s, '목 차')
    lines = [f"{(i + 1):02d}. {m['name']}  [{m['code'].upper()}]"
             for i, m in enumerate(menus)]
    parts = [{'text': line, 'size': 16, 'color': COLOR_DARK, 'space_after': 8}
             for line in lines]
    add_multipart_textbox(s, 0.6, 0.9, SLIDE_W - 1.2, SLIDE_H - 1.3, parts)
    return s


def build_menu_section_slide(prs, menu):
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    add_solid_rect(s, 0, 0, SLIDE_W, SLIDE_H,
                   fill=RGBColor(0xF4, 0xF6, 0xFA),
                   line=RGBColor(0xF4, 0xF6, 0xFA), line_w=0.5)
    add_solid_rect(s, 0, 3.0, SLIDE_W, 1.0,
                   fill=COLOR_TITLE_BG, line=COLOR_TITLE_BG, line_w=0.5)
    add_textbox(s, 0, 3.0, SLIDE_W, 0.6, menu['name'],
                size=36, bold=True, color=COLOR_WHITE,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 0, 3.55, SLIDE_W, 0.4, f"[{menu['code'].upper()}]",
                size=18, color=COLOR_SUBTITLE,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return s


def build_screen_slide(prs, *, title, image_path, geom, regions, desc_lines,
                       desc_x, desc_w):
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    add_title_bar(s, title)

    if image_path and image_path.exists():
        s.shapes.add_picture(
            str(image_path),
            Inches(geom.img_l), Inches(geom.img_t),
            width=Inches(geom.disp_w), height=Inches(geom.disp_h),
        )
    else:
        # 이미지 없음 placeholder
        add_solid_rect(s, geom.img_l, geom.img_t, geom.disp_w, geom.disp_h,
                       fill=RGBColor(0xEE, 0xF2, 0xF7),
                       line=COLOR_PANEL_LINE, line_w=1.0)
        add_textbox(s, geom.img_l, geom.img_t, geom.disp_w, geom.disp_h,
                    '(이미지 없음)', size=16, color=RGBColor(0x90, 0xA4, 0xAE),
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    if regions:
        add_region_labels(s, regions, geom)

    add_desc_panel(s, desc_x, desc_w, desc_lines)
    return s


# ── 메뉴별 시나리오 → regions / desc 합성 ──────────────────────
SCENARIOS = [
    {'file': '01-main.png', 'kind': 'main', 'title': '메인 화면'},
    {'file': '02-search-result.png', 'kind': 'search', 'title': '검색 결과'},
    {'file': '03-register-popup.png', 'kind': 'register', 'title': '등록 팝업'},
    {'file': '04-row-selected.png', 'kind': 'rowSelect', 'title': '행 선택'},
    {'file': '05-edit-popup.png', 'kind': 'edit', 'title': '수정 팝업'},
]


def bbox_to_px(bb):
    if not bb or bb.get('width', 0) <= 0:
        return None
    return [bb['x'], bb['y'],
            bb['x'] + bb['width'], bb['y'] + bb['height']]


def synth_regions_desc(kind, coords, ui):
    regions, desc = [], []
    if kind == 'main':
        c = coords.get('main', {})
        sb = bbox_to_px(c.get('search'))
        gb = bbox_to_px(c.get('grid'))
        if sb:
            regions.append({'px': sb, 'label': '① 검색 조건', 'color': COLOR_RED})
            desc.append({'h': '① 검색 조건', 'c': COLOR_RED})
            if ui['searchFields']:
                desc += ['· ' + f for f in ui['searchFields']]
            else:
                desc.append('· 화면 상단 조건을 입력 후 [검색] 클릭')
            desc.append('')
        if gb:
            regions.append({'px': gb, 'label': '② 결과 그리드', 'color': COLOR_GRAY})
            desc.append({'h': '② 결과 그리드', 'c': COLOR_GRAY})
            desc.append('· 검색 실행 전 빈 상태로 표시됩니다.')
            desc.append('')
        if ui['rules']:
            desc.append({'h': '업무 규칙', 'c': COLOR_NAVY})
            desc += ['· ' + r for r in ui['rules']]
    elif kind == 'search':
        c = coords.get('search', {})
        gb = bbox_to_px(c.get('grid'))
        if gb:
            regions.append({'px': gb, 'label': '결과 그리드', 'color': COLOR_GREEN})
        desc.append({'h': '검색 결과', 'c': COLOR_GREEN})
        desc.append('· 입력한 조건에 해당하는 데이터가 표시됩니다.')
        if ui['gridColumns']:
            desc.append('')
            desc.append({'h': '주요 컬럼', 'c': COLOR_NAVY})
            desc += ['· ' + col for col in ui['gridColumns']]
        desc.append('')
        desc.append({'h': '사용 방법', 'c': COLOR_NAVY})
        desc.append('1. 결과 행을 클릭하면 상세 정보가 표시됩니다.')
        desc.append('2. 페이지 하단에서 페이지 이동 / 페이지 크기를 변경할 수 있습니다.')
    elif kind in ('register', 'edit'):
        key = 'register' if kind == 'register' else 'edit'
        c = coords.get(key, {})
        pp = bbox_to_px(c.get('popup'))
        label = '등록 팝업' if kind == 'register' else '수정 팝업'
        if pp:
            regions.append({'px': pp, 'label': label, 'color': COLOR_NAVY})
        desc.append({'h': label, 'c': COLOR_NAVY})
        if kind == 'register':
            desc.append('· [추가] 버튼 클릭 시 표시됩니다.')
        else:
            desc.append('· 결과 그리드에서 행 선택 후 [수정] 클릭 시 표시됩니다.')
        desc.append('· 필수 항목(*)을 모두 입력한 후 [저장] 버튼을 클릭합니다.')
        desc.append('· [취소] 또는 [✕] 클릭 시 변경사항 없이 닫힙니다.')
        desc.append('')
        desc.append('⚠ 본 매뉴얼은 팝업 화면을 보여주기 위해 열기만 한 상태입니다.')
    elif kind == 'rowSelect':
        c = coords.get('rowSelect', {})
        gb = bbox_to_px(c.get('grid'))
        if gb:
            regions.append({'px': gb, 'label': '선택된 행', 'color': COLOR_BLUE})
        desc.append({'h': '행 선택', 'c': COLOR_BLUE})
        desc.append('· 그리드 행 클릭 시 해당 행이 강조 표시됩니다.')
        desc.append('· 행 선택 후 [수정] / [삭제] / [복사] 등의 기능을 실행할 수 있습니다.')
    return regions, desc


# ── 메인 ────────────────────────────────────────────────────────
def main():
    prs = Presentation(str(TEMPLATE))
    # 슬라이드 크기 보정 (안전상)
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # 템플릿의 예제 슬라이드는 모두 제거 (마스터/레이아웃/테마는 유지)
    remove_all_slides(prs)

    # 표지
    build_cover_slide(prs, CUSTOMER)
    # 목차
    build_toc_slide(prs, cfg['menus'])

    viewport = cfg.get('viewport', {'width': 1440, 'height': 900})
    img_col_w = img_col_w_for(viewport)
    desc_x = img_col_w
    desc_w = SLIDE_W - img_col_w

    # 메뉴별 슬라이드
    for menu in cfg['menus']:
        screen_dir = SCREENS_ROOT / menu['code']
        coords_path = screen_dir / 'coords.json'
        coords = (json.loads(coords_path.read_text(encoding='utf-8'))
                  if coords_path.exists() else {})
        ui = parse_ui_md(menu['code'])

        # 메뉴 섹션 표지
        build_menu_section_slide(prs, menu)

        # 시나리오별 슬라이드
        for sc in SCENARIOS:
            img_path = screen_dir / sc['file']
            if not img_path.exists():
                continue
            size = png_size(img_path)
            px_w, px_h = size if size else (viewport['width'], viewport['height'])
            geom = Geom(px_w, px_h, img_col_w)

            regions, desc = synth_regions_desc(sc['kind'], coords, ui)
            title = f"{menu['name']} - {sc['title']}  [{menu['code'].upper()}]"
            build_screen_slide(prs, title=title, image_path=img_path,
                               geom=geom, regions=regions,
                               desc_lines=desc, desc_x=desc_x, desc_w=desc_w)

    # 페이지 번호 (모든 슬라이드 마지막에 일괄)
    total = len(prs.slides)
    for i, sl in enumerate(prs.slides, start=1):
        add_page_number(sl, i, total)

    prs.save(str(OUT_FILE))

    print(f"[OK] PPTX 생성 완료")
    print(f"     {OUT_FILE}")
    print(f"     총 슬라이드 {total}장 (메뉴 {len(cfg['menus'])}개)")
    print(f"     템플릿     : {TEMPLATE}")


if __name__ == '__main__':
    main()
